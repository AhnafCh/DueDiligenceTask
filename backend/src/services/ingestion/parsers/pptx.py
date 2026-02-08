"""PPTX document parser using python-pptx."""
import logging
from pathlib import Path
from typing import List
from pptx import Presentation

from .base import BaseParser, ParsedContent

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Parser for Microsoft PowerPoint presentations."""
    
    def supports(self, file_extension: str) -> bool:
        """Check if file extension is .pptx"""
        return file_extension.lower() == ".pptx"
    
    def parse(self, file_path: Path) -> ParsedContent:
        """
        Parse PPTX document and extract text from slides.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            ParsedContent with extracted text and metadata
            
        Raises:
            ValueError: If PPTX is invalid
            IOError: If file cannot be read
        """
        try:
            prs = Presentation(file_path)
            
            slide_texts: List[str] = []
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = [f"=== Slide {slide_num} ==="]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                slide_content.append(row_text)
                
                # Extract notes if available
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_content.append(f"[Notes: {notes_text}]")
                    except Exception as e:
                        logger.debug(f"Could not extract notes from slide {slide_num}: {e}")
                
                # Add slide content if not empty
                if len(slide_content) > 1:
                    slide_texts.append("\n".join(slide_content))
                else:
                    slide_texts.append(f"=== Slide {slide_num} ===\n[Empty slide]")
            
            # Combine all slides
            full_text = "\n\n".join(slide_texts)
            cleaned_text = self.clean_text(full_text)
            
            # Each slide is treated as a page
            page_count = len(prs.slides)
            page_texts = slide_texts
            
            # Extract metadata
            metadata = self.extract_metadata(
                page_count=page_count,
                slide_count=len(prs.slides),
                file_size=file_path.stat().st_size,
                file_name=file_path.name
            )
            
            # Add presentation properties if available
            try:
                core_props = prs.core_properties
                if core_props.title:
                    metadata['title'] = core_props.title
                if core_props.author:
                    metadata['author'] = core_props.author
                if core_props.created:
                    metadata['created'] = str(core_props.created)
                if core_props.modified:
                    metadata['modified'] = str(core_props.modified)
            except Exception as e:
                logger.debug(f"Could not extract PPTX metadata: {e}")
            
            return ParsedContent(
                text=cleaned_text,
                page_count=page_count,
                metadata=metadata,
                page_texts=page_texts,
                bounding_boxes=None  # PPTX doesn't have bounding boxes
            )
            
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise IOError(f"Failed to parse PPTX: {e}")
